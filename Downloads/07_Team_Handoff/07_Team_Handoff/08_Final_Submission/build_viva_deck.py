"""Build the 5011CEM viva presentation deck (16:9, python-pptx).

Generates 5011CEM_Viva_Presentation.pptx in the same folder as this script.
All numbers come from 06_Working_Files/REPORT_FACTS_PACK.md (reproduced 2026-06-11).
Run:  python -X utf8 build_viva_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- paths
HERE = Path(__file__).resolve().parent
ROOT = HERE.parent  # 07_Team_Handoff root
EDA = ROOT / "04_Diagrams" / "EDA_Plots"
SYSD = ROOT / "04_Diagrams" / "System_Diagrams"
FIGS = ROOT / "05_GitHub_and_Code" / "code" / "reports" / "figures"
OUT = HERE / "5011CEM_Viva_Presentation.pptx"

# ---------------------------------------------------------------- style
DARK = RGBColor(0x1F, 0x38, 0x64)    # dark blue accent
ACCENT = RGBColor(0x2E, 0x74, 0xB5)  # mid blue
BODY = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xDC, 0xE6, 0xF1)    # light blue table banding
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SLIDE_TITLES = []


def new_slide(title=None):
    slide = prs.slides.add_slide(BLANK)
    if title is not None:
        # thin top strip
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.10))
        strip.fill.solid()
        strip.fill.fore_color.rgb = DARK
        strip.line.fill.background()
        # small accent block left of the title
        block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.34), Inches(0.14), Inches(0.52))
        block.fill.solid()
        block.fill.fore_color.rgb = ACCENT
        block.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.24), Inches(12.1), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = DARK
        SLIDE_TITLES.append(title)
        # footer page number
        n = len(prs.slides.__iter__.__self__._sldIdLst)  # current count
        ftb = slide.shapes.add_textbox(Inches(12.45), Inches(7.08), Inches(0.7), Inches(0.32))
        fp = ftb.text_frame.paragraphs[0]
        fr = fp.add_run()
        fr.text = str(n)
        fr.font.size = Pt(12)
        fr.font.name = FONT
        fr.font.color.rgb = MUTED
        fp.alignment = PP_ALIGN.RIGHT
    return slide


def add_bullets(slide, items, left, top, width, height, size=18, space=6):
    """items: list of str or (str, level) or (str, level, bold)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        bold = False
        level = 0
        if isinstance(item, tuple):
            if len(item) == 3:
                text, level, bold = item
            else:
                text, level = item
        else:
            text = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space)
        p.level = level
        r = p.add_run()
        r.text = ("- " if level == 0 else "  - ") + text
        r.font.size = Pt(size if level == 0 else max(size - 2, 14))
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = BODY
    return tb


def add_text(slide, text, left, top, width, height, size=16, bold=False,
             color=None, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = FONT
    r.font.color.rgb = color if color is not None else BODY
    return tb


def add_picture_fit(slide, path, left, top, max_w, max_h):
    """Insert picture scaled to fit inside the box, centred."""
    pic = slide.shapes.add_picture(str(path), left, top)
    scale = min(max_w / pic.width, max_h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(left + (max_w - pic.width) / 2)
    pic.top = int(top + (max_h - pic.height) / 2)
    return pic


def add_caption(slide, text, left, top, width):
    return add_text(slide, text, left, top, width, Inches(0.3), size=12,
                    color=MUTED, align=PP_ALIGN.CENTER, italic=True)


def add_table(slide, data, left, top, width, height, font_size=13,
              col_widths=None, all_left=False):
    rows, cols = len(data), len(data[0])
    gfx = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = gfx.table
    if col_widths:
        total = float(sum(col_widths))
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / total)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(font_size)
            r.font.name = FONT
            cell.fill.solid()
            if ri == 0:
                r.font.bold = True
                r.font.color.rgb = WHITE
                cell.fill.fore_color.rgb = DARK
            else:
                r.font.color.rgb = BODY
                cell.fill.fore_color.rgb = BAND if ri % 2 == 0 else WHITE
            if all_left or ci == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER
    return tbl


# ================================================================ slide 1: title
s = new_slide()  # no header bar on title slide
SLIDE_TITLES.append("Title")
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.7),
                          prs.slide_width, Inches(2.3))
band.fill.solid()
band.fill.fore_color.rgb = DARK
band.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Home Credit Loan-Default Prediction"
r.font.size = Pt(40)
r.font.bold = True
r.font.name = FONT
r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
p2.space_before = Pt(8)
r2 = p2.add_run()
r2.text = "5011CEM Big Data Programming Project  |  Viva Presentation"
r2.font.size = Pt(20)
r2.font.name = FONT
r2.font.color.rgb = RGBColor(0xCF, 0xDD, 0xF2)
add_text(s, "Bryan Tey Kai Yuan (P23015693)   |   Thong Wai Kit (P23015668)",
         Inches(0.8), Inches(4.45), Inches(11.7), Inches(0.5), size=20, bold=True, color=DARK)
add_text(s, "Lecturer: Ms. Vimala Doraisamy   |   April 2026 Semester",
         Inches(0.8), Inches(5.05), Inches(11.7), Inches(0.4), size=16, color=BODY)
add_text(s, "INTI International College Penang, School of Computing, "
            "in collaboration with Coventry University, UK",
         Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.4), size=14, color=MUTED)

# ================================================================ slide 2: problem
s = new_slide("Problem and Objectives")
add_bullets(s, [
    ("Problem: predict which loan applicants will face payment difficulty so the "
     "lender can manage credit risk before approval", 0, True),
    "Default is rare (8.07%), so the real challenge is finding the minority class "
    "without rejecting too many good customers",
    "Objectives: build an end-to-end pipeline (load, EDA, preprocess, train, "
    "evaluate, optimise, dashboard)",
    "Compare three classifiers (Logistic Regression, Random Forest, XGBoost) under "
    "class imbalance and justify the choice with evidence",
    "Deliver explainable outputs: statistical tests, feature importance, and an "
    "analyst-facing Streamlit dashboard",
], Inches(0.7), Inches(1.4), Inches(12.0), Inches(4.6), size=20, space=12)
add_text(s, "Success criteria: best model selected by F1 and ROC-AUC on a held-out "
            "test set, reproducible one-command pipeline, tested code.",
         Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.6), size=16, italic=True, color=ACCENT)

# ================================================================ slide 3: dataset
s = new_slide("Dataset: Home Credit Default Risk (Kaggle)")
add_bullets(s, [
    "application_train.csv: 307,511 rows x 122 columns (166 MB), one row per application",
    "TARGET = 1 means payment difficulty: 24,825 defaults = 8.07% (91.9% vs 8.1% imbalance)",
    "About 24% of all cells missing; sparse building-survey block mostly empty",
    "Quirk: DAYS_EMPLOYED sentinel 365243 marks pensioners / unemployed; replaced "
    "with NaN on load (18.0% of rows)",
    "Real, anonymised consumer-loan data from Home Credit Group (2018); no direct "
    "identifiers; meets the 100,000+ record requirement",
], Inches(0.6), Inches(1.35), Inches(6.7), Inches(5.0), size=18, space=10)
add_picture_fit(s, EDA / "eda_target.png", Inches(7.5), Inches(1.35), Inches(5.4), Inches(5.2))
add_caption(s, "Class balance: repaid vs default", Inches(7.5), Inches(6.7), Inches(5.4))

# ================================================================ slide 4: EDA
s = new_slide("EDA Highlights")
add_bullets(s, [
    "Default rate falls with education: Lower secondary about 11%, Academic degree lowest",
    "Male applicants default more (about 10.1%) than female (about 7.0%)",
    "Income and credit amounts are right-skewed; log transform used for inspection",
    "External credit scores (EXT_SOURCE_1/2/3) show the strongest negative link to default",
], Inches(0.6), Inches(1.3), Inches(12.2), Inches(1.7), size=17, space=4)
add_picture_fit(s, EDA / "eda_default_by_cat.png", Inches(0.5), Inches(3.0), Inches(6.2), Inches(3.9))
add_caption(s, "Default rate by category", Inches(0.5), Inches(6.95), Inches(6.2))
add_picture_fit(s, EDA / "eda_corr.png", Inches(6.9), Inches(3.0), Inches(6.0), Inches(3.9))
add_caption(s, "Correlation heatmap (numeric features)", Inches(6.9), Inches(6.95), Inches(6.0))

# ================================================================ slide 5: preprocessing
s = new_slide("Preprocessing Pipeline (src/preprocessing.py)")
add_bullets(s, [
    "clean_data: drop exact duplicates (0 found); keep missing rows for imputation",
    "engineer_features: 6 domain ratios, e.g. AGE_YEARS, CREDIT_INCOME_RATIO, "
    "ANNUITY_INCOME_RATIO, CREDIT_TERM",
    "Drop SK_ID_CURR and TARGET, columns with >= 60% missing, and FLAG_DOCUMENT_* flags",
    "ColumnTransformer fit on train only (no leakage): numeric = median impute + "
    "StandardScaler; categorical = mode impute + one-hot (handle_unknown='ignore')",
    "Stratified 80/20 split keeps the 8.07% rate: 246,008 train / 61,503 test, "
    "random_state=42",
], Inches(0.6), Inches(1.35), Inches(6.9), Inches(5.4), size=17, space=10)
add_picture_fit(s, SYSD / "flowchart_preprocessing.png", Inches(7.7), Inches(1.3), Inches(5.2), Inches(5.5))
add_caption(s, "Preprocessing flowchart", Inches(7.7), Inches(6.9), Inches(5.2))

# ================================================================ slide 6: statistics
s = new_slide("Statistical Evidence for Variable Selection")
add_text(s, "Chi-square: categorical vs TARGET (all p < 0.001)",
         Inches(0.6), Inches(1.25), Inches(6.0), Inches(0.35), size=16, bold=True, color=ACCENT)
add_table(s, [
    ["Variable", "Chi-square", "dof"],
    ["NAME_INCOME_TYPE", "1253.5", "7"],
    ["NAME_EDUCATION_TYPE", "1019.2", "4"],
    ["CODE_GENDER", "920.8", "2"],
    ["NAME_FAMILY_STATUS", "504.7", "5"],
    ["NAME_CONTRACT_TYPE", "293.2", "1"],
], Inches(0.6), Inches(1.7), Inches(5.9), Inches(2.7), font_size=14, col_widths=[3, 1.4, 0.8])
add_text(s, "Point-biserial: numeric vs TARGET (all p < 0.05)",
         Inches(7.0), Inches(1.25), Inches(6.0), Inches(0.35), size=16, bold=True, color=ACCENT)
add_table(s, [
    ["Variable", "r"],
    ["EXT_SOURCE_3", "-0.179"],
    ["EXT_SOURCE_2", "-0.161"],
    ["EXT_SOURCE_1", "-0.155"],
    ["AGE_YEARS", "-0.078"],
    ["YEARS_EMPLOYED", "-0.075"],
    ["AMT_CREDIT", "-0.030"],
], Inches(7.0), Inches(1.7), Inches(5.7), Inches(3.1), font_size=14, col_widths=[2.4, 1.2])
add_bullets(s, [
    "External credit scores are the strongest single predictors of default",
    "Younger and shorter-employed applicants default more often",
    "Raw income barely matters once affordability ratios are included (r = -0.004)",
], Inches(0.6), Inches(5.35), Inches(12.2), Inches(1.6), size=17, space=6)

# ================================================================ slide 7: models
s = new_slide("Models and Imbalance Strategy")
add_table(s, [
    ["Model", "Key settings", "Imbalance handling"],
    ["Logistic Regression", "liblinear, max_iter=500", "class_weight='balanced'"],
    ["Random Forest", "200 trees, depth 20, min_leaf 10", "class_weight='balanced_subsample'"],
    ["XGBoost (deployed)", "400 trees, depth 6, lr 0.1, hist, aucpr", "scale_pos_weight = neg/pos (about 11.4)"],
], Inches(0.6), Inches(1.35), Inches(12.1), Inches(1.8), font_size=14, col_widths=[2.2, 3.4, 3.4])
add_bullets(s, [
    ("Evidence-based choice (xgb_experiment.py on the same hold-out):", 0, True),
    ("Baseline spw=1 vs spw-only vs SMOTE 0.3 vs SMOTE 0.5 vs combined", 1),
    ("scale_pos_weight alone gave the best F1 at the 0.5 threshold", 1),
    ("SMOTE plus spw double-corrects: recall rises but precision is crushed", 1),
    "Rule adopted: exactly one imbalance correction, applied inside the sklearn "
    "Pipeline so preprocessing and resampling never touch the test fold",
    "All three models share the same preprocessor for a fair comparison",
], Inches(0.6), Inches(3.45), Inches(12.1), Inches(3.4), size=18, space=8)

# ================================================================ slide 8: results
s = new_slide("Hold-out Results (61,503 applications)")
add_table(s, [
    ["Model", "Accuracy", "Precision", "Recall", "F1", "ROC-AUC"],
    ["XGBoost", "0.759", "0.190", "0.612", "0.290", "0.763"],
    ["Random Forest", "0.861", "0.240", "0.332", "0.278", "0.742"],
    ["Logistic Regression", "0.691", "0.162", "0.676", "0.261", "0.749"],
], Inches(0.6), Inches(1.3), Inches(12.1), Inches(1.55), font_size=14,
    col_widths=[2.6, 1.2, 1.2, 1.2, 1.2, 1.2])
add_bullets(s, [
    "At 8% defaults accuracy misleads: rejecting nobody already scores 92%",
    "XGBoost balances recall (0.612) and precision best: top F1 and AUC, so it is deployed",
    "RF is most accurate but misses two thirds of defaulters; LR catches most "
    "defaulters but flags many good customers",
], Inches(0.6), Inches(3.0), Inches(12.1), Inches(1.3), size=16, space=4)
add_picture_fit(s, FIGS / "cm_xgboost.png", Inches(0.9), Inches(4.25), Inches(5.6), Inches(2.7))
add_caption(s, "XGBoost confusion matrix", Inches(0.9), Inches(7.0), Inches(5.6))
add_picture_fit(s, FIGS / "roc_xgboost.png", Inches(6.9), Inches(4.25), Inches(5.6), Inches(2.7))
add_caption(s, "XGBoost ROC curve (AUC 0.763)", Inches(6.9), Inches(7.0), Inches(5.6))

# ================================================================ slide 9: threshold
s = new_slide("Optimisation 1: Decision-Threshold Tuning")
add_bullets(s, [
    "Sweep the classification threshold on the hold-out set and pick the F1 maximum",
    ("XGBoost: F1 0.290 to 0.315 at t = 0.65 (precision 0.190 to 0.260, recall "
     "0.612 to 0.400)", 0, True),
    "Logistic Regression: F1 0.261 to 0.300 at t = 0.66; Random Forest: 0.278 to "
    "0.285 at t = 0.47",
    "The threshold is a business lever: trade recall against precision to match "
    "the lender's review capacity and cost of a missed default",
], Inches(0.6), Inches(1.35), Inches(5.9), Inches(5.2), size=18, space=12)
add_picture_fit(s, FIGS / "threshold_sweep.png", Inches(6.7), Inches(1.35), Inches(6.2), Inches(5.3))
add_caption(s, "F1 vs threshold for all three models", Inches(6.7), Inches(6.75), Inches(6.2))

# ================================================================ slide 10: complexity + calibration
s = new_slide("Optimisation 2: Complexity, Runtime, Calibration")
add_table(s, [
    ["Model", "Fit (s)", "Predict (s)", "Train complexity"],
    ["Logistic Regression", "95.0", "0.5", "O(n*d) per epoch, single-threaded"],
    ["Random Forest", "58.8", "1.2", "O(T*n*log n*sqrt(d)), T=200"],
    ["XGBoost", "22.8", "0.9", "O(T*L*n) with histogram binning, T=400"],
], Inches(0.6), Inches(1.3), Inches(7.0), Inches(1.9), font_size=13,
    col_widths=[2.2, 0.9, 1.1, 3.2])
add_bullets(s, [
    "XGBoost trains fastest despite the most trees: hist bins features once and "
    "parallelises; liblinear is single-threaded on about 240 one-hot columns",
    "Parquet cache cuts reloads of the 166 MB CSV several-fold",
    "Brier scores: RF 0.124, XGB 0.163, LR 0.202; all curves sit above the diagonal",
    "Weighted models over-score, so outputs are treated as ranking scores, not "
    "probabilities; the dashboard shows percentile risk bands instead",
], Inches(0.6), Inches(3.4), Inches(7.0), Inches(3.6), size=16, space=8)
add_picture_fit(s, FIGS / "calibration_curves.png", Inches(7.9), Inches(1.3), Inches(5.0), Inches(5.4))
add_caption(s, "Calibration curves (reliability diagram)", Inches(7.9), Inches(6.8), Inches(5.0))

# ================================================================ slide 11: importance
s = new_slide("Feature Importance: What Drives Default Risk")
add_bullets(s, [
    ("EXT_SOURCE_3, EXT_SOURCE_2, EXT_SOURCE_1 dominate for both RF and XGBoost", 0, True),
    "Next: CREDIT_TERM, AGE_YEARS / DAYS_BIRTH, DAYS_EMPLOYED",
    "XGBoost gain also ranks Higher education, CODE_GENDER and Revolving loans",
    "LR's largest coefficients sit on rare categories (Academic degree, Student, "
    "Pensioner income types), a known weakness of unregularised one-hot effects",
    "Consistent with Yang et al. (2025): EXT_SOURCE features dominate SHAP values "
    "on this dataset",
], Inches(0.6), Inches(1.35), Inches(5.9), Inches(5.2), size=18, space=12)
add_picture_fit(s, FIGS / "feature_importance.png", Inches(6.7), Inches(1.35), Inches(6.2), Inches(5.3))
add_caption(s, "Top features per model", Inches(6.7), Inches(6.75), Inches(6.2))

# ================================================================ slide 12: dashboard
s = new_slide("Dashboard Walkthrough (Streamlit, 3 pages)")
add_table(s, [
    ["Page", "What it shows", "Built by"],
    ["Overview", "KPIs + interactive EDA charts on a 50,000-row sample", "Bryan"],
    ["Model Comparison", "Metrics table, threshold-tuning table, confusion matrices, "
     "ROC, calibration and importance figures (reads reports/metrics.csv)", "Thong"],
    ["Live Prediction", "Applicant form, XGBoost score, percentile risk band "
     "Low / Moderate / High using sample score quantiles q50 / q80", "Bryan"],
], Inches(0.6), Inches(1.35), Inches(12.1), Inches(2.6), font_size=14,
    col_widths=[1.8, 6.4, 1.0])
add_bullets(s, [
    "Single source of truth: the dashboard reads the same metrics.csv and figures "
    "that run_pipeline.py writes; no numbers are typed in twice",
    "Risk bands avoid quoting inflated raw probabilities from the weighted model",
], Inches(0.6), Inches(4.2), Inches(12.1), Inches(1.2), size=17, space=6)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), Inches(5.6),
                         Inches(8.0), Inches(1.3))
box.fill.solid()
box.fill.fore_color.rgb = DARK
box.line.fill.background()
btf = box.text_frame
btf.word_wrap = True
bp = btf.paragraphs[0]
bp.alignment = PP_ALIGN.CENTER
br = bp.add_run()
br.text = "LIVE DEMO: streamlit run dashboard/app.py"
br.font.size = Pt(22)
br.font.bold = True
br.font.name = FONT
br.font.color.rgb = WHITE
bp2 = btf.add_paragraph()
bp2.alignment = PP_ALIGN.CENTER
br2 = bp2.add_run()
br2.text = "Overview -> Model Comparison -> Live Prediction (high-risk and low-risk applicant)"
br2.font.size = Pt(14)
br2.font.name = FONT
br2.font.color.rgb = RGBColor(0xCF, 0xDD, 0xF2)

# ================================================================ slide 13: architecture
s = new_slide("System Architecture and Data Flow")
add_bullets(s, [
    "Raw CSV (166 MB) -> data_loader (parquet cache + sentinel fix) -> "
    "preprocessing -> three model pipelines -> evaluation outputs",
    "run_pipeline.py orchestrates everything and writes reports/metrics.csv plus "
    "all figures automatically",
    "The dashboard is a pure consumer of those outputs; retraining never requires "
    "touching dashboard code",
], Inches(0.6), Inches(1.3), Inches(12.2), Inches(1.55), size=17, space=4)
add_picture_fit(s, SYSD / "dfd_level1.png", Inches(1.4), Inches(2.85), Inches(10.5), Inches(4.0))
add_caption(s, "Data flow diagram (level 1)", Inches(1.4), Inches(6.9), Inches(10.5))

# ================================================================ slide 14: code tour
s = new_slide("Code Tour: Each File and Its Role")
add_table(s, [
    ["File", "Role", "Led by"],
    ["src/data_loader.py", "Load CSV, parquet cache, fix DAYS_EMPLOYED sentinel", "Thong"],
    ["src/preprocessing.py", "Cleaning, 6 ratio features, ColumnTransformer, stratified split", "Thong"],
    ["src/eda.py", "EDA figures + chi-square / point-biserial test tables", "Bryan"],
    ["src/models/ (LR, RF)", "Logistic Regression and Random Forest pipelines", "Bryan"],
    ["src/models/xgboost_model.py", "XGBoost pipeline with scale_pos_weight", "Thong"],
    ["src/evaluation.py", "Metrics, confusion matrices, ROC curves, metrics.csv", "Thong"],
    ["run_pipeline.py", "One command: EDA -> train 3 models -> consolidated outputs", "Both"],
    ["xgb_experiment.py", "Imbalance evidence: scale_pos_weight vs SMOTE", "Thong"],
    ["optimise_models.py", "Threshold sweep, calibration, timings, importances", "Thong"],
    ["dashboard/app.py", "Streamlit app: Overview, Model Comparison, Live Prediction", "Both"],
    ["tests/ (3 files)", "20 pytest unit tests + Streamlit AppTest smoke test", "Both"],
], Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.4), font_size=12,
    col_widths=[2.6, 5.6, 0.9], all_left=True)

# ================================================================ slide 15: professional practices
s = new_slide("Professional Practices")
add_bullets(s, [
    ("Version control:", 0, True),
    ("GitHub private repo (gokuthong/5011cem), trunk-based workflow, 36 commits "
     "(18 per member), descriptive conventional-commit messages, 15 May - 25 Jun", 1),
    ("Collaboration: clear work split, each member reviews the other's commits", 1),
    ("Testing:", 0, True),
    ("20 pytest unit tests (preprocessing, optimisation) + Streamlit AppTest smoke "
     "test for the dashboard, all passing", 1),
    ("Ethics and legal:", 0, True),
    ("Anonymised open data; GDPR Article 22: right not to be subject to a solely "
     "automated decision, so the model supports analysts, it does not auto-reject", 1),
    ("Fairness: gender is predictive but using it raises discrimination concerns "
     "under equal-credit rules; flagged as an ethical limitation", 1),
], Inches(0.6), Inches(1.35), Inches(12.2), Inches(5.6), size=18, space=6)

# ================================================================ slide 16: limitations
s = new_slide("Limitations and Future Work")
add_bullets(s, [
    ("Limitations:", 0, True),
    ("Single-table model: 6 relational child tables (about 2.7 GB) not yet used", 1),
    ("Scores are not calibrated probabilities (Brier 0.163 for XGB); imbalance "
     "weighting inflates raw scores", 1),
    ("Precision is modest (0.260 at the tuned threshold): about 3 false alarms per "
     "caught defaulter", 1),
    ("Demographic features (gender, age) raise fairness concerns", 1),
    ("Future work:", 0, True),
    ("Aggregate bureau / previous-application features on SK_ID_CURR", 1),
    ("Add a calibration step (Platt / isotonic) and cost-based threshold per "
     "business loss matrix", 1),
    ("Fairness audit, drift monitoring, and scalable storage (parquet now; MongoDB "
     "discussed for the semi-structured multi-table extension)", 1),
], Inches(0.6), Inches(1.35), Inches(12.2), Inches(5.6), size=18, space=6)

# ================================================================ slide 17: reflection
s = new_slide("Critical Reflection")
add_text(s, "Bryan Tey Kai Yuan", Inches(0.6), Inches(1.3), Inches(5.9), Inches(0.4),
         size=18, bold=True, color=ACCENT)
add_bullets(s, [
    "Strength: EDA depth and the dashboard turned numbers into decisions",
    "Weakness: first judged models on accuracy; the 0.861-accuracy RF that misses "
    "two thirds of defaulters corrected that instinct",
    "Lesson: statistical evidence (chi-square, correlations) makes feature choices "
    "defensible, not just visual",
], Inches(0.6), Inches(1.75), Inches(5.9), Inches(3.4), size=16, space=8)
add_text(s, "Thong Wai Kit", Inches(6.9), Inches(1.3), Inches(5.9), Inches(0.4),
         size=18, bold=True, color=ACCENT)
add_bullets(s, [
    "Strength: leakage-free pipeline design and the controlled imbalance experiment",
    "Weakness: underestimated calibration; only the Brier analysis showed the raw "
    "scores could not be quoted as probabilities",
    "Lesson: measure, do not assume; timing runs reversed my expectation that more "
    "trees means slower training",
], Inches(6.9), Inches(1.75), Inches(5.9), Inches(3.4), size=16, space=8)
add_text(s, "On each other's work: Bryan's EDA shaped Thong's feature engineering; "
            "Thong's evaluation module exposed weaknesses in Bryan's baselines early. "
            "Reviewing every commit caught issues (and kept the merge history clean).",
         Inches(0.6), Inches(5.5), Inches(12.2), Inches(1.2), size=16, italic=True, color=BODY)

# ================================================================ slide 18: closing
s = new_slide("Conclusion and Q&A")
add_bullets(s, [
    "XGBoost with scale_pos_weight is the deployed model: F1 0.315 at t = 0.65, "
    "ROC-AUC 0.763 on 61,503 held-out applications",
    "External credit scores dominate; affordability ratios beat raw income",
    "Reproducible, tested, version-controlled pipeline with a live dashboard",
    "Decisions stay human-in-the-loop, consistent with GDPR Article 22",
], Inches(0.7), Inches(1.5), Inches(12.0), Inches(3.2), size=20, space=12)
qband = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.3),
                           prs.slide_width, Inches(1.4))
qband.fill.solid()
qband.fill.fore_color.rgb = DARK
qband.line.fill.background()
qtf = qband.text_frame
qp = qtf.paragraphs[0]
qp.alignment = PP_ALIGN.CENTER
qr = qp.add_run()
qr.text = "Thank you. Questions?"
qr.font.size = Pt(32)
qr.font.bold = True
qr.font.name = FONT
qr.font.color.rgb = WHITE

# ================================================================ save
prs.save(str(OUT))
print("Saved:", OUT)
print("Slides:", len(prs.slides.__iter__.__self__._sldIdLst))
